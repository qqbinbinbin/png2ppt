#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation


def check_pptx(
    pptx_path: Path,
    source: Path | None = None,
    large_media_bytes: int = 500_000,
    require_next_to_source: bool = False,
    require_editable: bool = False,
    min_shapes: int = 0,
    min_text_shapes: int = 0,
    max_pictures: int | None = None,
) -> dict:
    pptx_path = pptx_path.expanduser().resolve()
    failures: list[str] = []
    if not pptx_path.exists():
        return {"pptx": str(pptx_path), "pass": False, "failures": ["missing_pptx"]}

    if source is not None:
        source = source.expanduser().resolve()
        if require_next_to_source and pptx_path.parent != source.parent:
            failures.append("final_not_next_to_source")

    prs = Presentation(str(pptx_path))
    slide_metrics = []
    total_shapes = 0
    total_pictures = 0
    total_text = 0
    for idx, slide in enumerate(prs.slides, 1):
        shapes = len(slide.shapes)
        pictures = sum(1 for shape in slide.shapes if shape.shape_type == 13)
        text_shapes = sum(
            1
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
        )
        slide_metrics.append(
            {
                "slide": idx,
                "shapes": shapes,
                "pictures": pictures,
                "text_shapes": text_shapes,
            }
        )
        total_shapes += shapes
        total_pictures += pictures
        total_text += text_shapes
        if require_editable and shapes <= pictures and text_shapes == 0:
            failures.append(f"slide_{idx}_not_editable")

    if len(prs.slides) == 0:
        failures.append("no_slides")
    if total_shapes < min_shapes:
        failures.append("too_few_shapes")
    if total_text < min_text_shapes:
        failures.append("too_few_text_shapes")
    if max_pictures is not None and total_pictures > max_pictures:
        failures.append("too_many_pictures")
    if require_editable and total_shapes <= total_pictures:
        failures.append("deck_not_editable")

    media = []
    with ZipFile(pptx_path) as archive:
        for name in archive.namelist():
            if not name.startswith("ppt/media/"):
                continue
            size = archive.getinfo(name).file_size
            media.append({"name": name, "bytes": size})
    large_media = [item for item in media if item["bytes"] > large_media_bytes]

    result = {
        "pptx": str(pptx_path),
        "source": str(source) if source else None,
        "file_bytes": pptx_path.stat().st_size,
        "slides": len(prs.slides),
        "total_shapes": total_shapes,
        "total_pictures": total_pictures,
        "total_text_shapes": total_text,
        "slide_metrics": slide_metrics,
        "media_count": len(media),
        "large_media_threshold": large_media_bytes,
        "large_media": large_media,
        "failures": failures,
    }
    result["pass"] = not failures
    return result


def main() -> None:
    parser = argparse.ArgumentParser(description="Check final PPTX structure and packaging hygiene.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--source", type=Path)
    parser.add_argument("--out", type=Path)
    parser.add_argument("--large-media-bytes", type=int, default=500_000)
    parser.add_argument("--require-next-to-source", action="store_true")
    parser.add_argument("--require-editable", action="store_true")
    parser.add_argument("--min-shapes", type=int, default=0)
    parser.add_argument("--min-text-shapes", type=int, default=0)
    parser.add_argument("--max-pictures", type=int)
    parser.add_argument("--fail-on-large-media", action="store_true")
    args = parser.parse_args()

    result = check_pptx(
        args.pptx,
        source=args.source,
        large_media_bytes=args.large_media_bytes,
        require_next_to_source=args.require_next_to_source,
        require_editable=args.require_editable,
        min_shapes=args.min_shapes,
        min_text_shapes=args.min_text_shapes,
        max_pictures=args.max_pictures,
    )
    if args.fail_on_large_media and result["large_media"]:
        result["failures"].append("large_media")
        result["pass"] = False

    text = json.dumps(result, ensure_ascii=False, indent=2)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(text + "\n", encoding="utf-8")
    print(text)
    if not result["pass"]:
        raise SystemExit(2)


if __name__ == "__main__":
    main()
