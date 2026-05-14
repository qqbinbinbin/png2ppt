#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)


def image_size(path: Path) -> tuple[int, int]:
    from PIL import Image

    with Image.open(path) as im:
        return im.size


def add_image_slide(prs: Presentation, image: Path, fit: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    iw, ih = image_size(image)
    slide_ratio = SLIDE_W / SLIDE_H
    image_ratio = iw / ih

    if (fit == "contain" and image_ratio > slide_ratio) or (fit == "cover" and image_ratio < slide_ratio):
        w = SLIDE_W
        h = int(SLIDE_W / image_ratio)
    else:
        h = SLIDE_H
        w = int(SLIDE_H * image_ratio)

    left = int((SLIDE_W - w) / 2)
    top = int((SLIDE_H - h) / 2)
    slide.shapes.add_picture(str(image), left, top, width=w, height=h)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a PPTX deck from PNG files.")
    parser.add_argument("images", nargs="+", type=Path)
    parser.add_argument("--output", "-o", type=Path, required=True)
    parser.add_argument("--fit", choices=["contain", "cover"], default="contain")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    # Remove default first slide by creating a fresh deck with blank layout only in use.
    while len(prs.slides):
        rel_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[0]

    for source in args.images:
        image = source.expanduser().resolve()
        if not image.exists():
            raise FileNotFoundError(image)
        if image.suffix.lower() != ".png":
            raise ValueError(f"Only PNG input is supported: {image}")
        add_image_slide(prs, image, args.fit)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(args.output)


if __name__ == "__main__":
    main()
